"""将同目录下所有 .pptx 文件转换为 .md 笔记。

用法: python convert_pptx_to_md.py
输出: 每个 .pptx 生成同名 .md，放在同目录下。
"""
import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt

DIR = Path(__file__).parent

def shape_to_md(shape) -> str:
    """提取单个 shape 的文本，保留层级结构。"""
    if not shape.has_text_frame:
        return ""
    lines = []
    for para in shape.text_frame.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        # 根据缩进层级决定 markdown 列表深度
        level = para.level  # 0 = 顶层
        indent = "  " * level
        # 顶层且字号较大 → 当作标题
        if level == 0 and para.runs and para.runs[0].font.size and para.runs[0].font.size >= Pt(20):
            lines.append(f"## {text}")
        else:
            lines.append(f"{indent}- {text}")
    return "\n".join(lines)

def pptx_to_md(pptx_path: Path) -> str:
    """将单个 .pptx 转为 markdown 字符串。"""
    prs = Presentation(str(pptx_path))
    md_parts = [f"# {pptx_path.stem}\n"]

    for i, slide in enumerate(prs.slides, 1):
        slide_texts = []
        for shape in slide.shapes:
            txt = shape_to_md(shape)
            if txt:
                slide_texts.append(txt)

        if not slide_texts:
            continue

        md_parts.append(f"---\n\n### Slide {i}\n")
        md_parts.append("\n\n".join(slide_texts))

    return "\n\n".join(md_parts) + "\n"

def main():
    count = 0
    for f in sorted(DIR.glob("*.pptx")):
        if f.name.startswith("~$"):  # 跳过临时文件
            continue
        md_path = f.with_suffix(".md")
        md_content = pptx_to_md(f)
        md_path.write_text(md_content, encoding="utf-8")
        print(f"[OK] {f.name} -> {md_path.name}")
        count += 1
    print(f"\nDone: {count} files converted.")

if __name__ == "__main__":
    main()
